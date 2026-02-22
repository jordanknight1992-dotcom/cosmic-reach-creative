#!/usr/bin/env python3
"""
Cosmic Reach Creative — 36-Slide Brand Guidelines PPTX Generator
Generates a comprehensive brand guidelines deck with full visual styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Brand Tokens ───────────────────────────────────────────────────────────────
DEEP_SPACE     = RGBColor(0x0B, 0x11, 0x20)
SURFACE        = RGBColor(0x11, 0x18, 0x27)
SURFACE_ELEV   = RGBColor(0x1A, 0x23, 0x32)
STARLIGHT      = RGBColor(0xE8, 0xDF, 0xCF)
COPPER         = RGBColor(0xD4, 0xA5, 0x74)
SPARK_RED      = RGBColor(0xE0, 0x47, 0x47)
NAVY_ALT       = RGBColor(0x0B, 0x16, 0x28)
MUTED          = RGBColor(0x88, 0x96, 0xAB)
BORDER_CLR     = RGBColor(0x1E, 0x29, 0x3B)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
BLACK          = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Helpers ────────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 color=STARLIGHT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=STARLIGHT, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=0,
                  space_after=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p


def section_divider(prs, number, title, subtitle=""):
    """Creates a bold section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY_ALT)
    # Accent line
    add_rect(slide, Inches(1), Inches(2.8), Inches(0.8), Pt(4), SPARK_RED)
    # Number
    add_text_box(slide, Inches(1), Inches(3.1), Inches(3), Inches(0.8),
                 f"Section {number}", font_size=14, color=COPPER, bold=True)
    # Title
    add_text_box(slide, Inches(1), Inches(3.6), Inches(10), Inches(1.2),
                 title, font_size=44, color=STARLIGHT, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.8), Inches(8), Inches(0.8),
                     subtitle, font_size=18, color=MUTED)
    return slide


def content_slide(prs, label, title, body_lines, bg=DEEP_SPACE):
    """Standard content slide with label, title, and body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg)
    # Label
    add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
                 label.upper(), font_size=11, color=SPARK_RED, bold=True)
    # Title
    add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
                 title, font_size=36, color=STARLIGHT, bold=True)
    # Body
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10.5), Inches(4.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED
        p.font.name = "Calibri"
        p.space_after = Pt(10)
    return slide


def card_on_slide(slide, left, top, width, height, title, body, fill=SURFACE, border=BORDER_CLR):
    """Draw a card shape with title and body text."""
    card = add_rounded_rect(slide, left, top, width, height, fill, border)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3),
                 width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=STARLIGHT, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.9),
                 width - Inches(0.6), height - Inches(1.2),
                 body, font_size=12, color=MUTED)


# ─── Build Presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
# Decorative accent bar at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), SPARK_RED)
# Title
add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
             "Cosmic Reach Creative", font_size=54, color=STARLIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)
# Tagline
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
             "Strategy. Systems. Signal.", font_size=24, color=COPPER, bold=True,
             alignment=PP_ALIGN.CENTER)
# Subtitle
add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.6),
             "Brand Guidelines  |  2026", font_size=16, color=MUTED,
             alignment=PP_ALIGN.CENTER)
# Bottom accent
add_rect(slide, Inches(5.5), Inches(5.8), Inches(2.3), Pt(3), SPARK_RED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.5),
             "CONTENTS", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.8),
             "What's Inside", font_size=36, color=STARLIGHT, bold=True)

toc_items = [
    ("01", "Brand Foundation", "Mission, vision, beliefs, positioning"),
    ("02", "Visual Identity", "Logo, color, typography, imagery"),
    ("03", "Voice & Tone", "How we speak, write, and communicate"),
    ("04", "Service Architecture", "Creative, Parallax, Labs"),
    ("05", "Applications", "Digital, print, presentations, social"),
    ("06", "Brand in Action", "Process, governance, resources"),
]
for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(2.2) + Inches(0.8) * i
    add_text_box(slide, Inches(1), y, Inches(0.5), Inches(0.5),
                 num, font_size=14, color=SPARK_RED, bold=True)
    add_text_box(slide, Inches(1.7), y, Inches(3), Inches(0.5),
                 title, font_size=16, color=STARLIGHT, bold=True)
    add_text_box(slide, Inches(5), y, Inches(6), Inches(0.5),
                 desc, font_size=13, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Section Divider — Brand Foundation
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "01", "Brand Foundation", "Who we are. Why we exist. What we believe.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Brand Story / Origin
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Brand Story", "The Origin of Cosmic Reach", [
    "Jordan Knight built Cosmic Reach after leading marketing strategy and analytics inside a scaling life sciences organization, where clarity was operational necessity.",
    "",
    "Inside high growth environments, he built reporting systems, aligned cross functional teams, and translated performance data into executive level decisions.",
    "",
    "Cosmic Reach was launched to bring that same structured thinking to growing companies at the exact moment complexity outpaces structure.",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Mission Statement
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "MISSION", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Our Mission", font_size=36, color=STARLIGHT, bold=True)
# Large quote-style mission
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.5),
             "Systems design for companies navigating complexity.\nWe restore the signal when growth creates noise.",
             font_size=28, color=STARLIGHT, bold=True, alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1),
             "We work with growing companies at the exact moment complexity outpaces structure. We design the systems that restore signal, so your team can stop reacting and start navigating.",
             font_size=16, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Vision
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "VISION", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Our Vision", font_size=36, color=STARLIGHT, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.5),
             "Every growing company deserves the clarity\nto navigate complexity with confidence.",
             font_size=28, color=STARLIGHT, bold=True)
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(1.5),
             "We envision a world where scaling does not mean losing focus. Where data creates decisions, not confusion. Where teams move with precision because the systems beneath them hold.",
             font_size=16, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Brand Beliefs — VISUALLY STUNNING
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)

# Full-width gradient overlay at top (dark to transparent feel)
top_bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SPARK_RED)

# Large decorative geometric elements
# Left vertical accent strip
add_rect(slide, Inches(0), Inches(0.06), Inches(0.08), SLIDE_H, SURFACE_ELEV)
# Subtle grid pattern - horizontal lines
for i in range(8):
    y = Inches(0.9) + Inches(0.82) * i
    add_rect(slide, Inches(0.4), y, Inches(12.5), Pt(0.5), BORDER_CLR)

# Glowing orb decorative element (top right)
add_circle(slide, Inches(10.5), Inches(-0.3), Inches(2), SURFACE_ELEV)
add_circle(slide, Inches(10.8), Inches(0.0), Inches(1.4), SURFACE)
add_circle(slide, Inches(11.0), Inches(0.2), Inches(1.0), NAVY_ALT)

# Section label and title
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "BRAND BELIEFS", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(8), Inches(0.7),
             "What We Believe", font_size=40, color=STARLIGHT, bold=True)

# Beliefs in a dramatic staggered card layout with accent indicators
beliefs = [
    ("Clarity is a system, not a moment.",
     "Sustainable clarity comes from designed structure, not a single insight."),
    ("Signal over noise. Always.",
     "We filter the metrics that matter from the ones that just look busy."),
    ("Structure enables speed.",
     "The right system does not slow your team down. It removes the friction that does."),
    ("Complexity is not the enemy. Confusion is.",
     "Growing companies will always face complexity. We make sure it does not become chaos."),
    ("Every system should hold under pressure.",
     "We architect the simplest system that will hold, then ship, measure, and refine."),
]

for i, (belief, desc) in enumerate(beliefs):
    # Stagger left offset slightly for visual rhythm
    x_offset = Inches(0.8) + Inches(0.15) * (i % 3)
    y_pos = Inches(1.85) + Inches(1.05) * i
    card_w = Inches(11.5) - Inches(0.15) * (i % 3)

    # Card background
    card = add_rounded_rect(slide, x_offset, y_pos, card_w, Inches(0.9), SURFACE, BORDER_CLR)

    # Spark red accent pip on left edge
    add_rect(slide, x_offset, y_pos + Inches(0.15), Pt(5), Inches(0.6), SPARK_RED)

    # Belief text (bold, large)
    add_text_box(slide, x_offset + Inches(0.45), y_pos + Inches(0.08),
                 Inches(5.5), Inches(0.5),
                 belief, font_size=16, color=STARLIGHT, bold=True)

    # Description text (muted, right aligned area)
    add_text_box(slide, x_offset + Inches(6.2), y_pos + Inches(0.08),
                 Inches(5), Inches(0.5),
                 desc, font_size=12, color=MUTED)

    # Belief number indicator (copper circle)
    num_circle = add_circle(slide, x_offset + Inches(0.12), y_pos + Inches(0.22),
                            Inches(0.45), DEEP_SPACE, COPPER)
    # Number text
    add_text_box(slide, x_offset + Inches(0.12), y_pos + Inches(0.25),
                 Inches(0.45), Inches(0.4),
                 f"0{i+1}", font_size=11, color=COPPER, bold=True,
                 alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), SPARK_RED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Positioning
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "POSITIONING", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Brand Positioning", font_size=36, color=STARLIGHT, bold=True)

pos_items = [
    ("For", "Growing companies whose scale has outpaced their systems"),
    ("Who Need", "Structure that turns complexity into clarity"),
    ("We Are", "A systems design partner that restores the signal"),
    ("Unlike", "Traditional agencies that add more noise to an already crowded landscape"),
    ("We", "Design the systems, strategy, and structure that give your team clarity to move"),
]
for i, (label, desc) in enumerate(pos_items):
    y = Inches(2.5) + Inches(0.9) * i
    add_text_box(slide, Inches(1), y, Inches(1.5), Inches(0.5),
                 label, font_size=14, color=COPPER, bold=True)
    add_text_box(slide, Inches(2.8), y, Inches(9), Inches(0.5),
                 desc, font_size=16, color=STARLIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Key Messages
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "KEY MESSAGES", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Core Messaging Framework", font_size=36, color=STARLIGHT, bold=True)

messages = [
    ("Primary Headline", "Growth created noise. We restore the signal."),
    ("Elevator", "Systems design for companies navigating complexity."),
    ("Problem", "You scaled. Now everything is louder. More channels. More dashboards. More opinions."),
    ("Solution", "We design the systems that restore signal, so your team can stop reacting and start navigating."),
    ("CTA", "Request a Signal Scan"),
]
for i, (label, msg) in enumerate(messages):
    y = Inches(2.5) + Inches(0.9) * i
    card_on_slide(slide, Inches(1), y, Inches(11.3), Inches(0.75), label, msg)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Section Divider — Visual Identity
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "02", "Visual Identity", "Logo, color, typography, and imagery standards.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Logo System Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "LOGO SYSTEM", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Logo Architecture", font_size=36, color=STARLIGHT, bold=True)

add_text_box(slide, Inches(1), Inches(2.4), Inches(10), Inches(0.6),
             "The Cosmic Reach logo exists in a systematic family of 8 variants. Select the right version based on background, context, and medium.",
             font_size=16, color=MUTED)

# Logo variant grid
variants = [
    ("Full Logo — Dark RGB", "logo-dark-rgb.svg", "Primary. Dark backgrounds."),
    ("Full Logo — Dark Mono", "logo-dark-mono.svg", "Single-color dark contexts."),
    ("Full Logo — Light RGB", "logo-light-rgb.svg", "Light backgrounds."),
    ("Full Logo — Light Mono", "logo-light-mono.svg", "Single-color light contexts."),
    ("Icon — Dark RGB", "icon-dark-rgb.svg", "Favicons, app icons (dark)."),
    ("Icon — Dark Mono", "icon-dark-mono.svg", "Minimal dark contexts."),
    ("Icon — Light RGB", "icon-light-rgb.svg", "Favicons, app icons (light)."),
    ("Icon — Light Mono", "icon-light-mono.svg", "Minimal light contexts."),
]
for i, (name, file, usage) in enumerate(variants):
    col = i % 4
    row = i // 4
    x = Inches(1) + Inches(2.9) * col
    y = Inches(3.3) + Inches(1.8) * row
    card = add_rounded_rect(slide, x, y, Inches(2.7), Inches(1.6), SURFACE, BORDER_CLR)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.3), Inches(0.4),
                 name, font_size=11, color=STARLIGHT, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.3), Inches(0.3),
                 file, font_size=9, color=COPPER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.3), Inches(0.4),
                 usage, font_size=10, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Logo Usage — Clear Space, Sizing & Rules
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Logo Usage", "Clear Space, Sizing & Usage Rules", [
    "Clear Space: Maintain minimum clear space equal to the height of the 'C' in Cosmic on all sides.",
    "",
    "Minimum Size: Full logo 120px / 1in. Icon mark 32px / 0.5in.",
    "",
    "Do Not:",
    "  • Distort, stretch, rotate, or skew the logo",
    "  • Apply unapproved colors or add effects (shadows, glows)",
    "  • Place on low-contrast or busy backgrounds",
    "  • Crop, partially obscure, or rearrange logo elements",
    "  • Use the logo as a pattern or texture",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Color Palette — Primary
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "COLOR SYSTEM", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Primary Color Palette", font_size=36, color=STARLIGHT, bold=True)
add_text_box(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.5),
             "Dark-first design philosophy. Deep navy is dominant. Red is functional signal only.",
             font_size=16, color=MUTED)

colors = [
    ("Deep Space", "#0B1120", "Primary background", DEEP_SPACE),
    ("Surface", "#111827", "Cards, elevated areas", SURFACE),
    ("Starlight", "#E8DFCF", "Primary text", STARLIGHT),
    ("Copper", "#D4A574", "Secondary accent", COPPER),
    ("Spark Red", "#E04747", "CTAs, signals", SPARK_RED),
    ("Muted", "#8896AB", "Secondary text", MUTED),
]
for i, (name, hex_val, usage, rgb) in enumerate(colors):
    x = Inches(1) + Inches(2.0) * i
    # Color swatch
    swatch = add_rounded_rect(slide, x, Inches(3.0), Inches(1.7), Inches(1.7), rgb, BORDER_CLR)
    # Name
    text_color = DEEP_SPACE if name in ("Starlight", "Copper", "Spark Red") else STARLIGHT
    add_text_box(slide, x, Inches(4.9), Inches(1.7), Inches(0.4),
                 name, font_size=13, color=STARLIGHT, bold=True)
    add_text_box(slide, x, Inches(5.3), Inches(1.7), Inches(0.3),
                 hex_val, font_size=11, color=COPPER)
    add_text_box(slide, x, Inches(5.7), Inches(1.7), Inches(0.3),
                 usage, font_size=10, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Color Palette — Extended & Rules
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "COLOR RULES", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Color Usage Guidelines", font_size=36, color=STARLIGHT, bold=True)

# Extended colors
ext_colors = [
    ("Surface Elevated", "#1A2332", "Hover states", SURFACE_ELEV),
    ("Brand Navy Alt", "#0B1628", "Section variation", NAVY_ALT),
    ("Border", "#1E293B", "Dividers", BORDER_CLR),
]
for i, (name, hex_val, usage, rgb) in enumerate(ext_colors):
    x = Inches(1) + Inches(2.5) * i
    swatch = add_rounded_rect(slide, x, Inches(2.5), Inches(2.0), Inches(1.2), rgb, BORDER_CLR)
    add_text_box(slide, x, Inches(3.9), Inches(2.0), Inches(0.3),
                 f"{name}  {hex_val}", font_size=11, color=STARLIGHT, bold=True)
    add_text_box(slide, x, Inches(4.2), Inches(2.0), Inches(0.3),
                 usage, font_size=10, color=MUTED)

# Rules
rules = [
    "Dark-first: Deep navy is the dominant surface across all brand touchpoints",
    "Red is functional signal only: CTAs, alerts, active states. Never decorative",
    "No random gradients: Gradients are only used for subtle depth effects",
    "No neon or sci-fi visuals: The brand is grounded, structured, and warm",
    "WCAG 2.1 AA minimum contrast for all text pairings",
]
for i, rule in enumerate(rules):
    y = Inches(5.0) + Inches(0.4) * i
    add_rect(slide, Inches(1), y + Inches(0.12), Pt(8), Pt(8), SPARK_RED)
    add_text_box(slide, Inches(1.4), y, Inches(10), Inches(0.4),
                 rule, font_size=12, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Typography — Display Font
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "TYPOGRAPHY", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Display Font: Space Grotesk", font_size=36, color=STARLIGHT, bold=True)

add_text_box(slide, Inches(1), Inches(2.4), Inches(10), Inches(0.5),
             "Used for headings, buttons, labels, and navigation. A modern geometric sans-serif that carries authority without pretension.",
             font_size=16, color=MUTED)

# Weight samples
add_text_box(slide, Inches(1), Inches(3.3), Inches(5), Inches(0.6),
             "Semibold 600", font_size=14, color=COPPER, bold=True)
add_text_box(slide, Inches(1), Inches(3.8), Inches(10), Inches(0.6),
             "Subheadings, buttons, nav links", font_size=28, color=STARLIGHT)

add_text_box(slide, Inches(1), Inches(4.8), Inches(5), Inches(0.6),
             "Bold 700", font_size=14, color=COPPER, bold=True)
add_text_box(slide, Inches(1), Inches(5.3), Inches(10), Inches(0.6),
             "Primary Headings", font_size=36, color=STARLIGHT, bold=True)

# Sample alphabet
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "A B C D E F G H I J K L M N O P Q R S T U V W X Y Z  0 1 2 3 4 5 6 7 8 9",
             font_size=14, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: Typography — Body Font
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "TYPOGRAPHY", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Body Font: Inter", font_size=36, color=STARLIGHT, bold=True)

add_text_box(slide, Inches(1), Inches(2.4), Inches(10), Inches(0.5),
             "Used for body text, descriptions, and form content. Designed for screen readability with excellent legibility at all sizes.",
             font_size=16, color=MUTED)

add_text_box(slide, Inches(1), Inches(3.3), Inches(5), Inches(0.6),
             "Regular 400", font_size=14, color=COPPER, bold=True)
add_text_box(slide, Inches(1), Inches(3.8), Inches(10), Inches(0.8),
             "Body text for descriptions, paragraphs, and form content. Optimized for readability at 16px and above.",
             font_size=18, color=STARLIGHT)

add_text_box(slide, Inches(1), Inches(4.9), Inches(5), Inches(0.6),
             "Medium 500", font_size=14, color=COPPER, bold=True)
add_text_box(slide, Inches(1), Inches(5.4), Inches(10), Inches(0.8),
             "Emphasized body text for callouts and important information.",
             font_size=18, color=STARLIGHT, bold=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: Type Scale
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "TYPE SCALE", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Typographic Scale", font_size=36, color=STARLIGHT, bold=True)

scales = [
    ("H1 Hero", "72px / 36px mobile", "Bold 700", 36),
    ("H2 Section", "48px / 30px mobile", "Bold 700", 28),
    ("H3 Card", "20px", "Bold 700", 20),
    ("Body Large", "18px", "Regular 400", 18),
    ("Body", "16px", "Regular 400", 16),
    ("Small / Labels", "14px", "Semibold 600", 14),
    ("Caption", "12px", "Semibold 600", 12),
]
y_pos = Inches(2.5)
for name, size, weight, pt_size in scales:
    add_text_box(slide, Inches(1), y_pos, Inches(2.5), Inches(0.5),
                 name, font_size=12, color=COPPER, bold=True)
    add_text_box(slide, Inches(3.5), y_pos, Inches(2), Inches(0.5),
                 size, font_size=12, color=MUTED)
    add_text_box(slide, Inches(5.5), y_pos, Inches(2), Inches(0.5),
                 weight, font_size=12, color=MUTED)
    add_text_box(slide, Inches(8), y_pos - Inches(0.1), Inches(5), Inches(0.6),
                 "Signal", font_size=min(pt_size, 32), color=STARLIGHT, bold=True)
    y_pos += Inches(0.65)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19: Spacing & Grid
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "SPACING", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Spacing System & Grid", font_size=36, color=STARLIGHT, bold=True)
add_text_box(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.5),
             "Based on a 4px grid. All spacing uses multiples of 4px for consistent rhythm.",
             font_size=16, color=MUTED)

spacing = [
    ("xs", "4px"), ("sm", "8px"), ("md", "16px"), ("lg", "24px"),
    ("xl", "32px"), ("2xl", "48px"), ("3xl", "64px"), ("4xl", "96px"),
]
for i, (token, px) in enumerate(spacing):
    x = Inches(1) + Inches(1.4) * i
    bar_h = max(Inches(0.1), Inches(0.05 * (i + 1)))
    add_rect(slide, x, Inches(3.5), Inches(1.0), bar_h, SPARK_RED)
    add_text_box(slide, x, Inches(3.5) + bar_h + Inches(0.1), Inches(1.0), Inches(0.3),
                 token, font_size=12, color=STARLIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.5) + bar_h + Inches(0.4), Inches(1.0), Inches(0.3),
                 px, font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

# Grid info
grid_info = [
    ("Layout", "Max width: 1200px"),
    ("Desktop", "12 columns"),
    ("Tablet", "8 columns"),
    ("Mobile", "4 columns"),
    ("Padding", "20px mobile, 32px desktop"),
    ("Radius", "sm: 10px, md: 14px, lg: 18px"),
]
for i, (label, value) in enumerate(grid_info):
    col = i % 3
    row = i // 3
    x = Inches(1) + Inches(3.8) * col
    y = Inches(5.3) + Inches(0.7) * row
    add_text_box(slide, x, y, Inches(1.5), Inches(0.4),
                 label, font_size=12, color=COPPER, bold=True)
    add_text_box(slide, x + Inches(1.5), y, Inches(2.3), Inches(0.4),
                 value, font_size=12, color=STARLIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20: Motion & Animation
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Motion", "Motion & Animation", [
    "Motion Tokens:",
    "  • Fast: 180ms — Micro-interactions, hover states",
    "  • Base: 240ms — Standard transitions",
    "  • Slow: 300ms — Page-level transitions",
    "",
    "Easing: cubic-bezier(0.16, 1, 0.3, 1) — ease-cosmic",
    "",
    "All animations respect prefers-reduced-motion: reduce.",
    "",
    "Principles:",
    "  • Motion serves function, never decoration",
    "  • Keep transitions subtle and purposeful",
    "  • Entry animations use fade + slight translate",
    "  • Interactive feedback is immediate (< 180ms)",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21: Imagery & Photography
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Imagery", "Photography & Imagery", [
    "Photography style is grounded and intentional. No stock-photo energy.",
    "",
    "Guidelines:",
    "  • Dark, moody tones that complement the deep-space palette",
    "  • Natural lighting preferred over studio setups",
    "  • Images should feel authentic and unforced",
    "  • Overlays use gradient-to-deep-space for text legibility",
    "  • Hero images: 100vw, opacity-60 with gradient overlay",
    "",
    "Image Treatment:",
    "  • All images optimized through Next.js Image component",
    "  • Card images contained within rounded-lg borders",
    "  • Aspect ratios: Hero (full viewport), Featured (4:3), Portrait (3:4)",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22: Iconography & Chevron Motif
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Iconography", "Iconography & Chevron Motif", [
    "The chevron is the singular brand motif. It serves as a directional indicator on CTAs.",
    "",
    "Chevron Rules:",
    "  • Functional only. Never decorative",
    "  • SVG at 16px default size",
    "  • Inherits current text color",
    "  • Animates on hover: gap increases from 1 to 2 (subtle shift)",
    "",
    "General Icon Guidelines:",
    "  • Keep icon use minimal. The brand is text-driven",
    "  • When icons are needed, use simple line-weight SVGs",
    "  • Color inherits from context (starlight, spark-red, copper)",
    "  • No filled/solid icon styles. Outline only",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23: Section Divider — Voice & Tone
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "03", "Voice & Tone", "How we speak, write, and communicate the brand.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24: Brand Voice
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "BRAND VOICE", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "How We Sound", font_size=36, color=STARLIGHT, bold=True)

voice_attrs = [
    ("Confident", "We speak with authority earned from experience. Not loud. Not aggressive. Clear.", SPARK_RED),
    ("Structured", "Our communication has architecture. Ideas are organized, sequenced, and intentional.", COPPER),
    ("Cosmic", "We use celestial metaphors as a lens for clarity: signal, orbit, navigation, terrain.", STARLIGHT),
]
for i, (attr, desc, accent) in enumerate(voice_attrs):
    y = Inches(2.5) + Inches(1.5) * i
    card = add_rounded_rect(slide, Inches(1), y, Inches(11.3), Inches(1.3), SURFACE, BORDER_CLR)
    add_rect(slide, Inches(1), y + Inches(0.2), Pt(5), Inches(0.9), accent)
    add_text_box(slide, Inches(1.5), y + Inches(0.2), Inches(3), Inches(0.5),
                 attr, font_size=20, color=STARLIGHT, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.7), Inches(10), Inches(0.5),
                 desc, font_size=14, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25: Writing Rules
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Writing", "Writing Rules", [
    'Use "we" voice. Never "I." The brand is the collective, not the individual.',
    "",
    "No em dashes. Use commas, periods, or restructure sentences.",
    "",
    "Cosmic metaphors are welcome but disciplined. Signal, orbit, terrain, navigate. Not every sentence needs to reference space.",
    "",
    "No corporate cliches: leverage, synergy, disrupt, pivot, ideate.",
    "",
    "Short sentences. Let ideas breathe.",
    "",
    "Founder reference: Jordan Knight, Managing Partner, Cosmic Reach. Do not name previous employers.",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26: Tone Spectrum
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "TONE SPECTRUM", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "We Are / We Are Not", font_size=36, color=STARLIGHT, bold=True)

we_are = ["Confident", "Structured", "Direct", "Warm", "Expert", "Curious"]
we_are_not = ["Aggressive", "Rigid", "Blunt", "Casual", "Condescending", "Reckless"]

add_text_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(0.5),
             "We Are", font_size=20, color=COPPER, bold=True)
add_text_box(slide, Inches(7), Inches(2.5), Inches(5), Inches(0.5),
             "We Are Not", font_size=20, color=MUTED, bold=True)

for i in range(6):
    y = Inches(3.2) + Inches(0.6) * i
    add_text_box(slide, Inches(1.3), y, Inches(4), Inches(0.4),
                 we_are[i], font_size=16, color=STARLIGHT)
    # Divider dot
    add_circle(slide, Inches(6.3), y + Inches(0.1), Inches(0.15), SPARK_RED)
    add_text_box(slide, Inches(7.3), y, Inches(4), Inches(0.4),
                 we_are_not[i], font_size=16, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 27: Section Divider — Service Architecture
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "04", "Service Architecture", "Creative. Parallax. Labs. Three systems, one mission.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 28: Service Overview — Three Pillars
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "SERVICE PILLARS", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Three Systems. One Mission.", font_size=36, color=STARLIGHT, bold=True)

pillars = [
    ("Creative", "Strategy and systems consulting that turns scattered activity into focused, measurable signal.",
     "Brand strategy, campaign architecture, reporting frameworks",
     "Your team makes decisions from data, not guesswork."),
    ("Parallax", "A unified reporting surface that pulls GA4, Search Console, and LinkedIn into one clean view.",
     "Custom dashboards, automated reporting, data integration",
     "One source of truth your whole team trusts."),
    ("Labs", "Custom systems thinking and digital experimentation for problems that need a new approach.",
     "Custom tools, workflow systems, decision frameworks",
     "A solution designed for your specific terrain."),
]
for i, (name, desc, deliverables, outcome) in enumerate(pillars):
    x = Inches(0.8) + Inches(4.0) * i
    card = add_rounded_rect(slide, x, Inches(2.5), Inches(3.8), Inches(4.5), SURFACE, BORDER_CLR)
    add_rect(slide, x, Inches(2.5), Inches(3.8), Pt(4), SPARK_RED)
    add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(3.2), Inches(0.5),
                 name, font_size=22, color=STARLIGHT, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.2), Inches(1.2),
                 desc, font_size=12, color=MUTED)
    add_text_box(slide, x + Inches(0.3), Inches(4.9), Inches(3.2), Inches(0.3),
                 "Deliverables", font_size=10, color=COPPER, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(5.2), Inches(3.2), Inches(0.7),
                 deliverables, font_size=10, color=MUTED)
    add_text_box(slide, x + Inches(0.3), Inches(6.0), Inches(3.2), Inches(0.6),
                 outcome, font_size=11, color=COPPER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 29: Process — How We Work
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "PROCESS", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "How We Work", font_size=36, color=STARLIGHT, bold=True)

steps = [
    ("01", "Diagnose", "Map the signal, locate interference.",
     ["Signal Map", "Narrative Architecture", "Strategic Coordinates"]),
    ("02", "Design", "Architect the simplest system that will hold.",
     ["Identity Systems", "Digital Interfaces", "Motion Framework"]),
    ("03", "Deploy", "Ship the work, measure, refine.",
     ["Web Experience", "Platform Implementation", "Performance Tracking"]),
]
for i, (num, title, desc, deliverables) in enumerate(steps):
    x = Inches(0.8) + Inches(4.0) * i
    # Number circle
    circle = add_circle(slide, x + Inches(1.5), Inches(2.6), Inches(0.9), DEEP_SPACE, SPARK_RED)
    add_text_box(slide, x + Inches(1.5), Inches(2.7), Inches(0.9), Inches(0.7),
                 num, font_size=20, color=SPARK_RED, bold=True, alignment=PP_ALIGN.CENTER)
    # Title + desc
    add_text_box(slide, x + Inches(0.3), Inches(3.7), Inches(3.4), Inches(0.5),
                 title, font_size=24, color=STARLIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(4.3), Inches(3.4), Inches(0.6),
                 desc, font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)
    # Deliverables
    add_text_box(slide, x + Inches(0.3), Inches(5.1), Inches(3.4), Inches(0.3),
                 "Deliverables", font_size=10, color=COPPER, bold=True, alignment=PP_ALIGN.CENTER)
    for j, d in enumerate(deliverables):
        add_text_box(slide, x + Inches(0.3), Inches(5.5) + Inches(0.35) * j, Inches(3.4), Inches(0.3),
                     d, font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 30: Outcomes / Signals We Optimize For
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "OUTCOMES", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Signals We Optimize For", font_size=36, color=STARLIGHT, bold=True)

signals = [
    ("Decision Velocity", "Teams move from weekly bottlenecks to daily clarity on what matters."),
    ("Reporting Clarity", "One dashboard. One source of truth. No more spreadsheet archaeology."),
    ("Strategic Alignment", "Everyone from leadership to execution working from the same signal."),
    ("Operational Focus", "Less noise in the system. Fewer meetings about meetings."),
]
for i, (metric, desc) in enumerate(signals):
    col = i % 2
    row = i // 2
    x = Inches(1) + Inches(5.8) * col
    y = Inches(2.6) + Inches(2.0) * row
    card = add_rounded_rect(slide, x, y, Inches(5.5), Inches(1.7), SURFACE, BORDER_CLR)
    add_circle(slide, x + Inches(0.3), y + Inches(0.35), Inches(0.3), SPARK_RED)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.3), Inches(4.2), Inches(0.5),
                 metric, font_size=18, color=STARLIGHT, bold=True)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.9), Inches(4.2), Inches(0.6),
                 desc, font_size=13, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 31: Section Divider — Applications
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "05", "Applications", "How the brand lives across digital and physical touchpoints.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 32: Digital Applications — Website
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Digital", "Website & Digital Presence", [
    "Primary digital home: cosmicreachcreative.com",
    "",
    "Built with Next.js 16, React 19, Tailwind CSS v4. Performance-first architecture.",
    "",
    "Key Pages:",
    "  • Homepage — Hero, pillars, proof, email capture",
    "  • Work — Case studies and proof of concept",
    "  • Approach — Orbital mechanics (Diagnose, Design, Deploy)",
    "  • About — Founder story and company positioning",
    "  • Labs — Experimental tools and frameworks",
    "  • Parallax — Unified reporting dashboard",
    "  • Clarity Session — Consultation booking",
    "",
    "All pages follow the brand token system for consistent visual treatment.",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 33: Presentation Guidelines
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Applications", "Presentation Guidelines", [
    "All presentations follow the dark-first design system:",
    "",
    "  • Background: Deep Space (#0B1120) or Navy Alt (#0B1628)",
    "  • Text: Starlight (#E8DFCF) for headings, Muted (#8896AB) for body",
    "  • Accents: Spark Red for labels, Copper for metadata",
    "",
    "Slide Structure:",
    "  • Section label (red, uppercase, 11pt) at top",
    "  • Title (starlight, bold, 36pt) below label",
    "  • Body content with consistent spacing",
    "",
    "Use cards (Surface bg, Border stroke, rounded corners) for grouping content.",
    "Maintain generous whitespace. Let the content breathe.",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 34: Social Media & Email
# ═══════════════════════════════════════════════════════════════════════════════
content_slide(prs, "Applications", "Social Media & Email", [
    "LinkedIn: Primary social channel. Professional tone with cosmic vocabulary.",
    "",
    "Email Communications:",
    "  • From: jordan@cosmicreachcreative.com",
    "  • Signature includes: Jordan Knight, Managing Partner, Cosmic Reach",
    "  • Dark email templates matching website palette when possible",
    "",
    "Social Post Guidelines:",
    "  • Lead with insight, not promotion",
    "  • Short, structured posts. Use line breaks for rhythm",
    "  • Cosmic metaphors used sparingly in social context",
    "  • No hashtag spam. 2-3 relevant tags maximum",
    "  • Images follow the dark, moody photography style",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 35: Section Divider — Brand in Action
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "06", "Brand in Action", "Governance, resources, and how to keep the brand alive.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 36: Accessibility & Compliance
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(4), Inches(0.4),
             "ACCESSIBILITY", font_size=11, color=SPARK_RED, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.9),
             "Accessibility & Compliance", font_size=36, color=STARLIGHT, bold=True)

acc_items = [
    ("WCAG 2.1 AA", "Minimum standard for all text contrast pairings across digital touchpoints"),
    ("Focus States", "2px solid Spark Red with 2px offset on all interactive elements"),
    ("Keyboard Nav", "All interactive elements are fully keyboard accessible"),
    ("Reduced Motion", "All animations respect prefers-reduced-motion globally"),
    ("Color Independence", "Color never conveys information alone. Always paired with text or icons"),
    ("Semantic HTML", "Proper heading hierarchy, ARIA labels, and landmark regions"),
]
for i, (title, desc) in enumerate(acc_items):
    col = i % 2
    row = i // 2
    x = Inches(1) + Inches(5.8) * col
    y = Inches(2.5) + Inches(1.5) * row
    card = add_rounded_rect(slide, x, y, Inches(5.5), Inches(1.2), SURFACE, BORDER_CLR)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(4.9), Inches(0.4),
                 title, font_size=14, color=STARLIGHT, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(4.9), Inches(0.5),
                 desc, font_size=11, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 37: Closing Slide (with governance info)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_SPACE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), SPARK_RED)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "Cosmic Reach Creative", font_size=48, color=STARLIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8),
             "Strategy. Systems. Signal.", font_size=24, color=COPPER, bold=True,
             alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.6), Inches(2.3), Pt(3), SPARK_RED)

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
             "cosmicreachcreative.com  |  jordan@cosmicreachcreative.com", font_size=14, color=MUTED,
             alignment=PP_ALIGN.CENTER)

# Governance footer
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
             "Brand Owner: Jordan Knight, Managing Partner  |  Quarterly Review Cadence", font_size=11, color=COPPER,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.7),
             "Resources: BRAND_TOKENS.md  |  CONTENT_EDITING.md  |  /public/logos/ (8 SVG variants)  |  /public/images/",
             font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Pt(4), SPARK_RED)

# ─── Save ───────────────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "Cosmic_Reach_Brand_Guidelines_2026.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
